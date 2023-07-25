import collections
import collections.abc
from pptx import Presentation

prs = Presentation(r".\files\mail.pptx")
slide = prs.slides[0]

shape_list = slide.shapes
shape_idx = {}

for idx, value in enumerate(shape_list):
    print(idx, value.name)
    shape_idx[value.name] = idx

print(shape_idx)

