import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy
import openpyxl

def copyContentsToPPT(s_list, s_idx, shape_name, txt, fontsize):
    shape = s_list[s_idx[shape_name]]
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = txt

    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(fontsize)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0, 0, 0)

# 편지봉투 열기
prs = Presentation(r".\files\mail.pptx")
slide = prs.slides[0]
shape_list = slide.shapes
shape_idx = {}
source_slide = prs.slides[0]

for idx, value in enumerate(shape_list):
    shape_idx[value.name] = idx

# 주소록 열기
wb = openpyxl.load_workbook(r".\files\codefilex.xlsx")
ws = wb['codefilex']

for i in range(2, 100): # 첫번째 줄은 제목이 나타남
    name = f"{ws.cell(i, 2).value} {ws.cell(i, 3).value} 귀하"
    postcode = ws.cell(i, 6).value
    if postcode == None:
        postcode = ' '
    addr = ws.cell(i, 7).value

    copyContentsToPPT(shape_list, shape_idx, 'name', name, 18)
    copyContentsToPPT(shape_list, shape_idx, 'address', addr, 16)

    copyContentsToPPT(shape_list, shape_idx, 'pCode1', postcode[0:1], 14)
    copyContentsToPPT(shape_list, shape_idx, 'pCode2', postcode[1:2], 14)
    copyContentsToPPT(shape_list, shape_idx, 'pCode3', postcode[2:3], 14)
    copyContentsToPPT(shape_list, shape_idx, 'pCode4', postcode[3:4], 14)
    copyContentsToPPT(shape_list, shape_idx, 'pCode5', postcode[4:5], 14)

    ### 편지봉투 복사
    slide_layout = prs.slide_layouts[6]  # 빈 페이지
    copy_slide = prs.slides.add_slide(slide_layout)

    for shape1 in source_slide.shapes:
        el = shape1.element
        newel = copy.deepcopy(el)
        copy_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

prs.save(r".\files\envelops2.pptx")