import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy
import openpyxl

# 편지봉투 열기
prs = Presentation(r".\files\mail.pptx")
slide = prs.slides[0]
shape_list = slide.shapes
shape_idx = {}
source_slide = prs.slides[0]

for idx, value in enumerate(shape_list):
    shape_idx[value.name] = idx

# 주소록 열기
wb = openpyxl.load_workbook(r".\files\codefilex.xlsx")
ws = wb['codefilex']

for i in range(2, 100): # 첫번째 줄은 제목이 나타남
    name = f"{ws.cell(i, 2).value} {ws.cell(i, 3).value} 귀하"
    postcode = ws.cell(i, 6).value
    addr = ws.cell(i, 7).value

    # print(name, postcode, addr)
    ### 받는 사람
    shape = shape_list[shape_idx['name']]
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = name

    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(18)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0, 0, 0)

    ### 주소
    shape = shape_list[shape_idx['address']]
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = addr

    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(16)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0, 0, 0)

    ### 우편번호
    shape = shape_list[shape_idx['pCode1']]
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = postcode
    except:
        ...

    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(14)
    font.bold = True
    font.italic = None
    font.color.rgb = RGBColor(0, 0, 0)

    if i < 99:
        ### 편지봉투 복사
        slide_layout = prs.slide_layouts[6]  # 빈 페이지
        copy_slide = prs.slides.add_slide(slide_layout)

        for shape1 in source_slide.shapes:
            el = shape1.element
            newel = copy.deepcopy(el)
            copy_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

prs.save(r".\files\envelops.pptx")