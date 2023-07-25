import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation(r".\files\mail.pptx")
slide = prs.slides[0]

shape_list = slide.shapes
shape_idx = {}

for idx, value in enumerate(shape_list):
    shape_idx[value.name] = idx

name = "name"
str = "소 혜 민 귀하"
fontsize = 24
fontbold = True

shape = shape_list[shape_idx[name]]
text_frame = shape.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = str

font = run.font
font.name = '맑은 고딕'
font.size = Pt(fontsize)
font.bold = fontbold
font.italic = None
font.color.rgb = RGBColor(0, 0, 0)

prs.save(r".\files\mail_newname.pptx")