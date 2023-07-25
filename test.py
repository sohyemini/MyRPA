import collections
import collections.abc
from pptx import Presentation

prs = Presentation(r".\files\mail.pptx")
slide = prs.slides[0]

shape_list = slide.shapes
shape_idx = {}

for idx, value in enumerate(shape_list):
    shape_idx[value.name] = idx
print(shape_idx)

idx = 0
for value in shape_list:
    shape_idx[value.name] = idx
    idx += 1
print(shape_idx)