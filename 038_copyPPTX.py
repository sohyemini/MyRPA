import collections
import collections.abc
from pptx import Presentation
import copy

prs = Presentation(r".\files\mail.pptx")
source_slide =prs.slides[0]

slide_layout = prs.slide_layouts[6] #빈 페이지
copy_slide = prs.slides.add_slide(slide_layout)

for shape in source_slide.shapes:
    el = shape.element
    newel = copy.deepcopy(el)
    copy_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

prs.save(r".\files\copymail.pptx")